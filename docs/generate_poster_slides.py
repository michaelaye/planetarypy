"""Generate 3 versions of a single-slide PlanetaryPy feature showcase poster.

Each slide is a 16:9 widescreen poster with large fonts,
readable from the back of a conference room.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

OUTDIR = Path(__file__).parent
LOGO = OUTDIR / "plpy_logo.png"

# ── Font Sizes ──────────────────────────────────────────────────
# Tweak these dicts to test different size combinations per version.

V1_SIZES = {
    "title": 54,           # "planetarypy" heading
    "subtitle": 24,        # tagline
    "pip": 20,             # pip install line
    "code": 17,            # code example lines
    "code_small": 15,      # slightly smaller code (multi-line imports)
    "comment": 15,         # inline comments
    "stats": 18,           # bold stat lines in summary box
    "footer": 16,          # bottom credits
    "spacer": 8,           # blank line gap
    "spacer_small": 6,     # smaller blank gap
}

V2_SIZES = {
    "title": 60,           # "planetary\npy" heading
    "subtitle": 22,        # tagline
    "pip": 18,             # pip install line
    "code": 16,            # code example lines
    "comment": 14,         # inline comments
    "footer_left": 14,     # left column footer
}

V3_SIZES = {
    "title": 72,           # giant centered title
    "pip": 24,             # pip install line
    "code": 17,            # standard code lines
    "code_large": 18,      # full-width strip code
    "import": 16,          # import statement lines
    "comment": 14,         # inline comments
    "footer": 16,          # bottom credits
}

# Colors
BLACK = RGBColor(0x1A, 0x1A, 0x2E)
DARK_BG = RGBColor(0x16, 0x21, 0x3E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
ACCENT_BLUE = RGBColor(0x42, 0xAF, 0xFA)
ACCENT_GREEN = RGBColor(0xA6, 0xE2, 0x2E)
ACCENT_ORANGE = RGBColor(0xFD, 0x97, 0x1F)
ACCENT_PINK = RGBColor(0xF9, 0x26, 0x72)
ACCENT_YELLOW = RGBColor(0xE6, 0xDB, 0x74)
CODE_BG = RGBColor(0x27, 0x2B, 0x33)
BOX_BG = RGBColor(0x1E, 0x29, 0x3B)
SUBTLE = RGBColor(0x88, 0x88, 0x88)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, font_name="Arial", alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_code_box(slide, left, top, width, height, code_lines, accent_color=ACCENT_BLUE,
                 title=None, title_color=None, font_size=14):
    """Add a rounded box with code inside."""
    # Background box
    shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = BOX_BG
    shape.line.color.rgb = accent_color
    shape.line.width = Pt(2)

    y_offset = top + 0.15
    if title:
        add_textbox(slide, left + 0.2, y_offset, width - 0.4, 0.4,
                    title, font_size=font_size + 6, color=title_color or accent_color,
                    bold=True)
        y_offset += 0.45

    for line in code_lines:
        text = line["text"]
        clr = line.get("color", LIGHT_GRAY)
        sz = line.get("size", font_size)
        b = line.get("bold", False)
        fn = line.get("font", "JetBrains Mono")
        add_textbox(slide, left + 0.25, y_offset, width - 0.5, 0.35,
                    text, font_size=sz, color=clr, bold=b, font_name=fn)
        y_offset += sz / 72 * 1.5  # rough line spacing
    return shape


# ════════════════════════════════════════════════════════════════
# VERSION 1: "Dark Grid" — 2x3 grid of feature boxes
# ════════════════════════════════════════════════════════════════
def make_v1():
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, BLACK)

    S = V1_SIZES
    # Title bar
    add_textbox(slide, 0.5, 0.2, 8, 0.9, "planetarypy",
                font_size=S["title"], color=WHITE, bold=True, font_name="JetBrains Mono")
    add_textbox(slide, 0.5, 1.0, 10, 0.5,
                "Python Tools for Planetary Science Data Access",
                font_size=S["subtitle"], color=ACCENT_BLUE, bold=False)

    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(13.5), Inches(0.2),
                                  Inches(2.2), Inches(1.76))

    # pip install
    add_textbox(slide, 0.5, 1.55, 6, 0.4, "pip install planetarypy",
                font_size=S["pip"], color=ACCENT_GREEN, font_name="JetBrains Mono")

    # 2x3 grid of feature boxes
    col_w = 4.8
    row_h_top = 3.0
    row_h_bot = 2.8
    gap = 0.3
    x_start = 0.5
    y_row1 = 2.1
    y_row2 = y_row1 + row_h_top + gap

    # Box 1: PDS Index Retrieval
    add_code_box(slide, x_start, y_row1, col_w, row_h_top,
                 title="PDS Index Retrieval", title_color=ACCENT_BLUE,
                 code_lines=[
                     {"text": "from planetarypy import pds", "color": ACCENT_YELLOW, "size": S["code"]},
                     {"text": "", "size": S["spacer"]},
                     {"text": 'df = pds.get_index("mro.ctx.edr")', "color": WHITE, "size": S["code"]},
                     {"text": "", "size": S["spacer"]},
                     {"text": "# 164,103 products as DataFrame", "color": SUBTLE, "size": S["comment"]},
                     {"text": "# Auto-downloads, caches as Parquet", "color": SUBTLE, "size": S["comment"]},
                     {"text": "# Auto-updates when stale", "color": SUBTLE, "size": S["comment"]},
                 ], accent_color=ACCENT_BLUE, font_size=S["code"])

    # Box 2: Metadata Lookup
    add_code_box(slide, x_start + col_w + gap, y_row1, col_w, row_h_top,
                 title="Metadata Lookup", title_color=ACCENT_GREEN,
                 code_lines=[
                     {"text": "df = pds.get_index('cassini.iss.index')", "color": WHITE, "size": S["code"]},
                     {"text": "", "size": S["spacer"]},
                     {"text": 'cols = ["PRODUCT_ID", "TARGET_NAME",', "color": ACCENT_YELLOW, "size": S["code_small"]},
                     {"text": '        "IMAGE_TIME", "FILTER_NAME"]', "color": ACCENT_YELLOW, "size": S["code_small"]},
                     {"text": "", "size": S["spacer"]},
                     {"text": "df[cols].query(", "color": WHITE, "size": S["code"]},
                     {"text": '    "TARGET_NAME == \'ENCELADUS\'")', "color": WHITE, "size": S["code"]},
                 ], accent_color=ACCENT_GREEN, font_size=S["code"])

    # Box 3: SPICE Kernels
    add_code_box(slide, x_start + 2 * (col_w + gap), y_row1, col_w, row_h_top,
                 title="Archived SPICE Kernels", title_color=ACCENT_ORANGE,
                 code_lines=[
                     {"text": "from planetarypy.spice \\", "color": ACCENT_YELLOW, "size": S["code_small"]},
                     {"text": "    import archived_kernels as ak", "color": ACCENT_YELLOW, "size": S["code_small"]},
                     {"text": "", "size": S["spacer"]},
                     {"text": 'mk = ak.get_metakernel_and_files(', "color": WHITE, "size": S["code"]},
                     {"text": '    "mro",', "color": WHITE, "size": S["code"]},
                     {"text": '    start="2024-01-01",', "color": WHITE, "size": S["code"]},
                     {"text": '    stop="2024-01-31")', "color": WHITE, "size": S["code"]},
                 ], accent_color=ACCENT_ORANGE, font_size=S["code"])

    # Box 4: PDS Catalog
    add_code_box(slide, x_start, y_row2, col_w, row_h_bot,
                 title="PDS Catalog: 65 Missions", title_color=ACCENT_PINK,
                 code_lines=[
                     {"text": "from planetarypy.catalog import \\", "color": ACCENT_YELLOW, "size": S["code_small"]},
                     {"text": "    list_missions, list_products", "color": ACCENT_YELLOW, "size": S["code_small"]},
                     {"text": "", "size": S["spacer"]},
                     {"text": "list_missions()", "color": WHITE, "size": S["code"]},
                     {"text": "# ['apollo','cassini','dawn','galileo',", "color": SUBTLE, "size": S["comment"]},
                     {"text": "#  'juno','lro','mro','voyager',...]", "color": SUBTLE, "size": S["comment"]},
                 ], accent_color=ACCENT_PINK, font_size=S["code"])

    # Box 5: Product Download
    add_code_box(slide, x_start + col_w + gap, y_row2, col_w, row_h_bot,
                 title="Direct Product Download", title_color=ACCENT_YELLOW,
                 code_lines=[
                     {"text": "from planetarypy.catalog \\", "color": ACCENT_YELLOW, "size": S["code_small"]},
                     {"text": "    import fetch_product", "color": ACCENT_YELLOW, "size": S["code_small"]},
                     {"text": "", "size": S["spacer"]},
                     {"text": 'fetch_product("mro.ctx.edr",', "color": WHITE, "size": S["code"]},
                     {"text": '  "P02_001916_2221_XI_42N027W")', "color": WHITE, "size": S["code"]},
                     {"text": "# → downloads .IMG to local cache", "color": SUBTLE, "size": S["comment"]},
                 ], accent_color=ACCENT_YELLOW, font_size=S["code"])

    # Box 6: Direct Data Access Stats
    add_code_box(slide, x_start + 2 * (col_w + gap), y_row2, col_w, row_h_bot,
                 title="Direct Data Access", title_color=ACCENT_BLUE,
                 code_lines=[
                     {"text": "58 product types", "color": WHITE, "size": S["stats"], "font": "Arial", "bold": True},
                     {"text": "29 instruments", "color": WHITE, "size": S["stats"], "font": "Arial", "bold": True},
                     {"text": "15 missions", "color": WHITE, "size": S["stats"], "font": "Arial", "bold": True},
                     {"text": "", "size": S["spacer_small"]},
                     {"text": "90+ PDS indexes", "color": SUBTLE, "size": S["comment"], "font": "Arial"},
                     {"text": "37 SPICE missions", "color": SUBTLE, "size": S["comment"], "font": "Arial"},
                 ], accent_color=ACCENT_BLUE, font_size=S["code"])

    # Footer
    add_textbox(slide, 0.5, 8.5, 15, 0.4,
                "michaelaye.github.io/planetarypy  ·  github.com/planetarypy  ·  BSD-3 License",
                font_size=S["footer"], color=SUBTLE, alignment=PP_ALIGN.CENTER)

    prs.save(str(OUTDIR / "planetarypy_poster_v1_dark_grid.pptx"))
    print("Saved v1: dark_grid")


# ════════════════════════════════════════════════════════════════
# VERSION 2: "Horizontal Flow" — left title column, right features
# ════════════════════════════════════════════════════════════════
def make_v2():
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    S = V2_SIZES
    # Left column: title + branding (4 inches wide)
    # Title background stripe
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(4.5), Inches(9))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x0D, 0x15, 0x2A)
    shape.line.fill.background()

    add_textbox(slide, 0.4, 0.8, 3.8, 1.2, "planetary\npy",
                font_size=S["title"], color=WHITE, bold=True, font_name="JetBrains Mono")

    add_textbox(slide, 0.4, 2.8, 3.8, 0.8,
                "Python Tools for\nPlanetary Science\nData Access",
                font_size=S["subtitle"], color=ACCENT_BLUE)

    add_textbox(slide, 0.4, 4.3, 3.8, 0.5, "pip install planetarypy",
                font_size=S["pip"], color=ACCENT_GREEN, font_name="JetBrains Mono")

    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(1.0), Inches(5.2),
                                  Inches(2.5), Inches(2.0))

    add_textbox(slide, 0.3, 8.2, 4, 0.4,
                "github.com/planetarypy",
                font_size=S["footer_left"], color=SUBTLE, alignment=PP_ALIGN.CENTER)

    # Right side: 4 feature boxes stacked
    rx = 5.0
    rw = 10.5
    bh = 1.85
    gap = 0.2

    y = 0.3
    add_code_box(slide, rx, y, rw, bh,
                 title="PDS Index Retrieval — 90+ Indexes", title_color=ACCENT_BLUE,
                 code_lines=[
                     {"text": 'from planetarypy import pds', "color": ACCENT_YELLOW, "size": S["code"]},
                     {"text": 'df = pds.get_index("mro.ctx.edr")          # 164,103 products → DataFrame', "color": WHITE, "size": S["code"]},
                     {"text": 'df = pds.get_index("cassini.iss.index")    # auto-download, Parquet cache, auto-update', "color": WHITE, "size": S["code"]},
                 ], accent_color=ACCENT_BLUE, font_size=S["code"])

    y += bh + gap
    add_code_box(slide, rx, y, rw, bh,
                 title="Metadata Search & Filtering", title_color=ACCENT_GREEN,
                 code_lines=[
                     {"text": 'ctx = pds.get_index("mro.ctx.edr")', "color": WHITE, "size": S["code"]},
                     {"text": 'hits = ctx.query("CENTER_LATITUDE > 80 & SCALED_PIXEL_WIDTH < 6")', "color": WHITE, "size": S["code"]},
                     {"text": '# Filter 164K images by location, resolution, time — standard pandas', "color": SUBTLE, "size": S["comment"]},
                 ], accent_color=ACCENT_GREEN, font_size=S["code"])

    y += bh + gap
    add_code_box(slide, rx, y, rw, bh,
                 title="Archived SPICE Kernels — 37 Missions", title_color=ACCENT_ORANGE,
                 code_lines=[
                     {"text": 'from planetarypy.spice import archived_kernels as ak', "color": ACCENT_YELLOW, "size": S["code"]},
                     {"text": 'mk = ak.get_metakernel_and_files("mro", start="2024-01-01", stop="2024-01-31")', "color": WHITE, "size": S["code"]},
                     {"text": '# Date-filtered kernel subsets via NAIF — returns loadable metakernel path', "color": SUBTLE, "size": S["comment"]},
                 ], accent_color=ACCENT_ORANGE, font_size=S["code"])

    y += bh + gap
    add_code_box(slide, rx, y, rw, bh + 0.4,
                 title="PDS Catalog — 65 Missions, 2042 Product Types", title_color=ACCENT_PINK,
                 code_lines=[
                     {"text": 'from planetarypy.catalog import list_missions, fetch_product', "color": ACCENT_YELLOW, "size": S["code"]},
                     {"text": 'list_missions()                  # 65 missions across the entire PDS archive', "color": WHITE, "size": S["code"]},
                     {"text": 'list_products("mro.ctx")         # ["edr"]  — discover available products', "color": WHITE, "size": S["code"]},
                     {"text": 'fetch_product("mro.ctx.edr", "P02_001916_2221_XI_42N027W")  # → local path', "color": WHITE, "size": S["code"]},
                 ], accent_color=ACCENT_PINK, font_size=S["code"])

    prs.save(str(OUTDIR / "planetarypy_poster_v2_horizontal.pptx"))
    print("Saved v2: horizontal")


# ════════════════════════════════════════════════════════════════
# VERSION 3: "Bold Minimal" — big title, 3 wide feature strips
# ════════════════════════════════════════════════════════════════
def make_v3():
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, RGBColor(0x0A, 0x0A, 0x0A))

    S = V3_SIZES
    # Giant centered title
    add_textbox(slide, 0, 0.3, 16, 1.2, "planetarypy",
                font_size=S["title"], color=WHITE, bold=True,
                font_name="JetBrains Mono", alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 0, 1.4, 16, 0.5,
                "pip install planetarypy",
                font_size=S["pip"], color=ACCENT_GREEN,
                font_name="JetBrains Mono", alignment=PP_ALIGN.CENTER)

    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(13.8), Inches(0.15),
                                  Inches(1.8), Inches(1.44))

    # Three wide horizontal strips
    strip_h = 1.95
    gap = 0.25
    y = 2.2
    margin = 0.6

    # Strip 1: Data Access (left: indexes, right: catalog)
    add_code_box(slide, margin, y, 7.2, strip_h,
                 title="Index Retrieval", title_color=ACCENT_BLUE,
                 code_lines=[
                     {"text": 'from planetarypy import pds', "color": ACCENT_YELLOW, "size": S["code"]},
                     {"text": 'df = pds.get_index("mro.ctx.edr")', "color": WHITE, "size": S["code"]},
                     {"text": "# 164K products → DataFrame, cached", "color": SUBTLE, "size": S["comment"]},
                 ], accent_color=ACCENT_BLUE, font_size=S["code"])

    add_code_box(slide, margin + 7.5, y, 7.3, strip_h,
                 title="Metadata Filtering", title_color=ACCENT_GREEN,
                 code_lines=[
                     {"text": 'ctx = pds.get_index("mro.ctx.edr")', "color": WHITE, "size": S["code"]},
                     {"text": 'ctx.query("CENTER_LATITUDE > 80")', "color": WHITE, "size": S["code"]},
                     {"text": "# Any column, standard pandas", "color": SUBTLE, "size": S["comment"]},
                 ], accent_color=ACCENT_GREEN, font_size=S["code"])

    y += strip_h + gap

    # Strip 2: SPICE + Catalog
    add_code_box(slide, margin, y, 7.2, strip_h,
                 title="SPICE Kernels — 37 Missions", title_color=ACCENT_ORANGE,
                 code_lines=[
                     {"text": "from planetarypy.spice \\", "color": ACCENT_YELLOW, "size": S["import"]},
                     {"text": "    import archived_kernels as ak", "color": ACCENT_YELLOW, "size": S["import"]},
                     {"text": 'ak.get_metakernel_and_files("mro",', "color": WHITE, "size": S["code"]},
                     {"text": '    start="2024-01", stop="2024-02")', "color": WHITE, "size": S["code"]},
                 ], accent_color=ACCENT_ORANGE, font_size=S["code"])

    add_code_box(slide, margin + 7.5, y, 7.3, strip_h,
                 title="PDS Catalog — 65 Missions", title_color=ACCENT_PINK,
                 code_lines=[
                     {"text": "from planetarypy.catalog import \\", "color": ACCENT_YELLOW, "size": S["import"]},
                     {"text": "    list_missions, list_products", "color": ACCENT_YELLOW, "size": S["import"]},
                     {"text": 'list_missions()    # 65 missions', "color": WHITE, "size": S["code"]},
                     {"text": 'list_products("cassini.iss")', "color": WHITE, "size": S["code"]},
                 ], accent_color=ACCENT_PINK, font_size=S["code"])

    y += strip_h + gap

    # Strip 3: Full-width product download
    add_code_box(slide, margin, y, 14.8, strip_h,
                 title="Direct Product Download", title_color=ACCENT_YELLOW,
                 code_lines=[
                     {"text": 'from planetarypy.catalog import fetch_product', "color": ACCENT_YELLOW, "size": S["code_large"]},
                     {"text": 'path = fetch_product("mro.ctx.edr", "P02_001916_2221_XI_42N027W")  # → downloads & caches locally', "color": WHITE, "size": S["code_large"]},
                     {"text": '# Resolves URLs via PDS indexes · DuckDB-backed catalog · 2042 product types across entire PDS', "color": SUBTLE, "size": S["comment"]},
                 ], accent_color=ACCENT_YELLOW, font_size=S["code_large"])

    # Footer
    add_textbox(slide, 0, 8.4, 16, 0.4,
                "github.com/planetarypy/planetarypy  ·  planetarypy.org  ·  BSD-3 License  ·  90+ PDS indexes  ·  37 SPICE missions",
                font_size=S["footer"], color=SUBTLE, alignment=PP_ALIGN.CENTER)

    prs.save(str(OUTDIR / "planetarypy_poster_v3_bold_minimal.pptx"))
    print("Saved v3: bold_minimal")


if __name__ == "__main__":
    make_v1()
    # make_v2()
    # make_v3()
    print("\nAll 3 versions saved to docs/")
